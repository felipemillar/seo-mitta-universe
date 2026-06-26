#!/usr/bin/env python3
"""
Convierte los 3 documentos del cliente a Markdown:
1. Marco Teórico Renting MITTA GO.pptx
2. MATERIAL LOP.pptx
3. Términos y Condiciones Rent a Car.docx

Destino: 01_Estrategia_y_Planes/ (base de conocimientos del proyecto)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from docx import Document

BASE = "/Users/fmillar/Proyectos_Desarrollo/INMedios_Mitta"
DEST = os.path.join(BASE, "InMedios_Mitta", "01_Estrategia_y_Planes")


def pptx_to_md(pptx_path, md_path, title):
    """Convert a PowerPoint file to Markdown."""
    prs = Presentation(pptx_path)
    lines = [f"# {title}\n"]
    lines.append(f"> Fuente: `{os.path.basename(pptx_path)}` — Documento oficial del cliente Mitta\n")
    lines.append("---\n")

    for i, slide in enumerate(prs.slides, 1):
        slide_title = None
        slide_body_parts = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if not text:
                        continue

                    # First non-empty text in slide = slide title
                    if slide_title is None and shape.shape_id == slide.shapes[0].shape_id:
                        slide_title = text
                    else:
                        slide_body_parts.append(text)

            if shape.has_table:
                table = shape.table
                headers = []
                for cell in table.rows[0].cells:
                    headers.append(cell.text.strip())

                slide_body_parts.append("")
                slide_body_parts.append("| " + " | ".join(headers) + " |")
                slide_body_parts.append("| " + " | ".join(["---"] * len(headers)) + " |")

                for row in table.rows[1:]:
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    slide_body_parts.append("| " + " | ".join(cells) + " |")
                slide_body_parts.append("")

        if slide_title:
            lines.append(f"\n## Slide {i}: {slide_title}\n")
        else:
            lines.append(f"\n## Slide {i}\n")

        for part in slide_body_parts:
            if part == "":
                lines.append("")
            else:
                lines.append(f"{part}")
        lines.append("")

    with open(md_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))

    print(f"  ✅ {os.path.basename(pptx_path)} → {os.path.basename(md_path)} ({len(prs.slides)} slides)")


def docx_to_md(docx_path, md_path, title):
    """Convert a Word document to Markdown."""
    doc = Document(docx_path)
    lines = [f"# {title}\n"]
    lines.append(f"> Fuente: `{os.path.basename(docx_path)}` — Documento oficial del cliente Mitta\n")
    lines.append("---\n")

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        style_name = para.style.name.lower() if para.style else ""

        if "heading 1" in style_name or "título 1" in style_name:
            lines.append(f"\n## {text}\n")
        elif "heading 2" in style_name or "título 2" in style_name:
            lines.append(f"\n### {text}\n")
        elif "heading 3" in style_name or "título 3" in style_name:
            lines.append(f"\n#### {text}\n")
        elif "list" in style_name or text.startswith(("•", "-", "·", "▪")):
            clean = text.lstrip("•-·▪ ")
            lines.append(f"- {clean}")
        else:
            lines.append(f"\n{text}")

    # Also extract tables
    for table in doc.tables:
        lines.append("")
        headers = [cell.text.strip().replace("\n", " ") for cell in table.rows[0].cells]
        lines.append("| " + " | ".join(headers) + " |")
        lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
        for row in table.rows[1:]:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            lines.append("| " + " | ".join(cells) + " |")
        lines.append("")

    with open(md_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))

    para_count = len([p for p in doc.paragraphs if p.text.strip()])
    table_count = len(doc.tables)
    print(f"  ✅ {os.path.basename(docx_path)} → {os.path.basename(md_path)} ({para_count} párrafos, {table_count} tablas)")


if __name__ == "__main__":
    print("🔄 Convirtiendo documentos del cliente a Markdown...\n")

    # 1. Marco Teórico Renting MITTA GO
    pptx_to_md(
        os.path.join(BASE, "Marco Teórico Renting MITTA GO.pptx"),
        os.path.join(DEST, "Marco_Teorico_Renting_MittaGO.md"),
        "Marco Teórico: Renting MITTA GO"
    )

    # 2. MATERIAL LOP (Leasing Operativo)
    pptx_to_md(
        os.path.join(BASE, "MATERIAL LOP.pptx"),
        os.path.join(DEST, "Material_Leasing_Operativo_LOP.md"),
        "Material: Leasing Operativo (LOP) — Mitta"
    )

    # 3. Términos y Condiciones Rent a Car
    docx_to_md(
        os.path.join(BASE, "Términos y Condiciones Rent a Car.docx"),
        os.path.join(DEST, "Terminos_y_Condiciones_Rent_a_Car.md"),
        "Términos y Condiciones Generales de Arriendo — Mitta Rent a Car"
    )

    print(f"\n📁 Archivos guardados en: {DEST}")
    print("🏷️  Listos para consulta del agente y del equipo.")
